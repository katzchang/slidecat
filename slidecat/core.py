"""Core functionality for splitting and merging PowerPoint files."""

import copy
from pathlib import Path
from typing import List, Optional
from pptx import Presentation


def split_presentation(input_file: Path, output_dir: Path, chunk_size: Optional[int] = None) -> List[Path]:
    """
    Split a PowerPoint file into individual slides or chunks of slides.

    Args:
        input_file: Path to the input PPTX file
        output_dir: Directory to save split files
        chunk_size: Number of slides per chunk. If None, split into individual slides.

    Returns:
        List of paths to the created files
    """
    if not input_file.exists():
        raise FileNotFoundError(f"Input file not found: {input_file}")

    # Create output directory if it doesn't exist
    output_dir.mkdir(parents=True, exist_ok=True)

    # Load the presentation to get slide count
    prs = Presentation(str(input_file))
    total_slides = len(prs.slides)

    if total_slides == 0:
        raise ValueError(f"No slides found in {input_file}")

    if chunk_size is not None and chunk_size < 1:
        raise ValueError(f"Chunk size must be at least 1, got {chunk_size}")

    created_files = []
    stem = input_file.stem

    # Default to 1 slide per file if chunk_size is not specified
    if chunk_size is None:
        chunk_size = 1

    # Split into chunks
    chunk_num = 1
    for start_idx in range(0, total_slides, chunk_size):
        end_idx = min(start_idx + chunk_size, total_slides)

        # Load a fresh copy for each chunk
        new_prs = Presentation(str(input_file))

        # Determine which slides to keep (0-indexed)
        slides_to_keep = set(range(start_idx, end_idx))
        slides_to_remove = [i for i in range(len(new_prs.slides)) if i not in slides_to_keep]

        # Remove slides in reverse order to avoid index issues
        for slide_idx in reversed(slides_to_remove):
            rId = new_prs.slides._sldIdLst[slide_idx].rId
            new_prs.part.drop_rel(rId)
            del new_prs.slides._sldIdLst[slide_idx]

        # Save the new presentation
        if chunk_size == 1:
            # Single slide: use original naming
            output_file = output_dir / f"{stem}_slide_{start_idx + 1:03d}.pptx"
        else:
            # Multiple slides: show range
            output_file = output_dir / f"{stem}_slides_{start_idx + 1:03d}-{end_idx:03d}.pptx"

        new_prs.save(str(output_file))
        created_files.append(output_file)
        chunk_num += 1

    return created_files


def merge_presentations(input_files: List[Path], output_file: Path) -> Path:
    """
    Merge multiple PowerPoint files into one.

    Args:
        input_files: List of paths to input PPTX files
        output_file: Path to the output merged file

    Returns:
        Path to the created merged file
    """
    if not input_files:
        raise ValueError("No input files provided")

    # Check all input files exist
    for file in input_files:
        if not file.exists():
            raise FileNotFoundError(f"Input file not found: {file}")

    # Start with the first presentation
    merged_prs = Presentation(str(input_files[0]))

    # Add slides from remaining presentations
    for input_file in input_files[1:]:
        prs = Presentation(str(input_file))

        for slide in prs.slides:
            # Use blank layout (usually index 6, but we'll try to find it)
            # If that fails, use the first available layout
            try:
                blank_layout = merged_prs.slide_layouts[6]
            except IndexError:
                blank_layout = merged_prs.slide_layouts[0]

            new_slide = merged_prs.slides.add_slide(blank_layout)

            # Copy all shapes from the original slide using deepcopy
            for shape in slide.shapes:
                el = shape.element
                newel = copy.deepcopy(el)
                new_slide.shapes._spTree.append(newel)

            # Copy slide background if present
            if slide.background:
                try:
                    new_slide.background.fill.solid()
                    new_slide.background.fill.fore_color.rgb = slide.background.fill.fore_color.rgb
                except:
                    # Ignore background copy errors
                    pass

    # Create output directory if needed
    output_file.parent.mkdir(parents=True, exist_ok=True)

    # Save merged presentation
    merged_prs.save(str(output_file))

    return output_file


def extract_slides(
    input_file: Path,
    output_file: Path,
    start: int,
    end: Optional[int] = None
) -> Path:
    """
    Extract a range of slides from a PowerPoint file.

    Args:
        input_file: Path to the input PPTX file
        output_file: Path to the output file
        start: Starting slide number (1-indexed)
        end: Ending slide number (1-indexed, inclusive). If None, extract to the end.

    Returns:
        Path to the created file
    """
    if not input_file.exists():
        raise FileNotFoundError(f"Input file not found: {input_file}")

    # Load the presentation
    prs = Presentation(str(input_file))
    total_slides = len(prs.slides)

    if total_slides == 0:
        raise ValueError(f"No slides found in {input_file}")

    # Validate slide numbers
    if start < 1 or start > total_slides:
        raise ValueError(f"Start slide {start} is out of range (1-{total_slides})")

    if end is None:
        end = total_slides

    if end < start or end > total_slides:
        raise ValueError(f"End slide {end} is out of range ({start}-{total_slides})")

    # Create new presentation from the input file
    new_prs = Presentation(str(input_file))

    # Determine which slides to keep (convert to 0-indexed)
    slides_to_keep = set(range(start - 1, end))
    slides_to_remove = [i for i in range(len(new_prs.slides)) if i not in slides_to_keep]

    # Remove slides in reverse order to avoid index issues
    for slide_idx in reversed(slides_to_remove):
        rId = new_prs.slides._sldIdLst[slide_idx].rId
        new_prs.part.drop_rel(rId)
        del new_prs.slides._sldIdLst[slide_idx]

    # Create output directory if needed
    output_file.parent.mkdir(parents=True, exist_ok=True)

    # Save the new presentation
    new_prs.save(str(output_file))

    return output_file


def verify_presentation(input_file: Path) -> dict:
    """
    Verify a PowerPoint file and return diagnostic information.

    Args:
        input_file: Path to the PPTX file to verify

    Returns:
        Dictionary with verification results
    """
    result = {
        "valid": False,
        "error": None,
        "slides": 0,
        "slide_details": [],
        "layouts": 0,
        "masters": 0,
    }

    if not input_file.exists():
        result["error"] = f"File not found: {input_file}"
        return result

    try:
        # Try to load the presentation
        prs = Presentation(str(input_file))
        result["slides"] = len(prs.slides)
        result["layouts"] = len(prs.slide_layouts)
        result["masters"] = len(prs.slide_masters)

        # Check each slide
        for idx, slide in enumerate(prs.slides, start=1):
            slide_info = {
                "number": idx,
                "shapes": 0,
                "has_title": False,
                "error": None,
            }

            try:
                slide_info["shapes"] = len(slide.shapes)

                # Check if slide has a title
                if slide.shapes.title:
                    slide_info["has_title"] = True

            except Exception as e:
                slide_info["error"] = str(e)

            result["slide_details"].append(slide_info)

        # If we got here, the file is valid
        result["valid"] = True

    except Exception as e:
        result["error"] = str(e)

    return result
